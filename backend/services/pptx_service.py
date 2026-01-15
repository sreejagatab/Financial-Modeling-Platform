"""PowerPoint export service for financial presentations."""

import io
from dataclasses import dataclass, field
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData


@dataclass
class SlideConfig:
    """Configuration for a presentation slide."""
    title: str
    subtitle: str = ""
    slide_type: str = "content"  # title, content, metrics, table, chart, comparison
    content: Any = None
    notes: str = ""


@dataclass
class PresentationConfig:
    """Configuration for PowerPoint generation."""
    title: str = "Financial Analysis"
    subtitle: str = ""
    company_name: str = ""
    prepared_by: str = ""
    date: str = field(default_factory=lambda: datetime.now().strftime("%B %d, %Y"))
    primary_color: Tuple[int, int, int] = (26, 51, 128)  # RGB
    accent_color: Tuple[int, int, int] = (51, 153, 76)
    footer_text: str = "Confidential"


class PowerPointExportService:
    """Service for generating PowerPoint presentations from financial data."""

    def __init__(self, config: Optional[PresentationConfig] = None):
        self.config = config or PresentationConfig()
        self.prs = Presentation()
        self._setup_slide_masters()

    def _setup_slide_masters(self):
        """Set up slide dimensions (16:9 widescreen)."""
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def generate_presentation(
        self,
        slides: List[SlideConfig],
        output_path: Optional[str] = None,
    ) -> bytes:
        """Generate a PowerPoint presentation.

        Args:
            slides: List of slide configurations
            output_path: Optional path to save the file

        Returns:
            PPTX content as bytes
        """
        self.prs = Presentation()
        self._setup_slide_masters()

        # Title slide
        self._add_title_slide()

        # Content slides
        for slide_config in slides:
            self._add_slide(slide_config)

        # Save to buffer
        buffer = io.BytesIO()
        self.prs.save(buffer)
        pptx_content = buffer.getvalue()
        buffer.close()

        if output_path:
            with open(output_path, 'wb') as f:
                f.write(pptx_content)

        return pptx_content

    def _add_title_slide(self):
        """Add the title slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Background accent
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*self.config.primary_color)
        shape.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.8),
            Inches(12.333), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.config.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Subtitle
        if self.config.subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.9),
                Inches(12.333), Inches(0.5)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = self.config.subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(220, 220, 220)

        # Metadata
        meta_y = 3.5
        if self.config.company_name:
            self._add_meta_line(slide, "Company:", self.config.company_name, meta_y)
            meta_y += 0.5
        if self.config.prepared_by:
            self._add_meta_line(slide, "Prepared By:", self.config.prepared_by, meta_y)
            meta_y += 0.5
        self._add_meta_line(slide, "Date:", self.config.date, meta_y)

    def _add_meta_line(self, slide, label: str, value: str, y_pos: float):
        """Add a metadata line to the slide."""
        label_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_pos),
            Inches(1.5), Inches(0.4)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(128, 128, 128)

        value_box = slide.shapes.add_textbox(
            Inches(2), Inches(y_pos),
            Inches(6), Inches(0.4)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(14)
        p.font.bold = True

    def _add_slide(self, config: SlideConfig):
        """Add a content slide based on configuration."""
        if config.slide_type == "metrics":
            self._add_metrics_slide(config)
        elif config.slide_type == "table":
            self._add_table_slide(config)
        elif config.slide_type == "chart":
            self._add_chart_slide(config)
        elif config.slide_type == "comparison":
            self._add_comparison_slide(config)
        else:
            self._add_content_slide(config)

    def _add_content_slide(self, config: SlideConfig):
        """Add a basic content slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        self._add_slide_header(slide, config.title)

        if config.content:
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5),
                Inches(12.333), Inches(5.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            if isinstance(config.content, str):
                p = tf.paragraphs[0]
                p.text = config.content
                p.font.size = Pt(18)
            elif isinstance(config.content, list):
                for i, item in enumerate(config.content):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(16)
                    p.space_after = Pt(12)

        if config.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = config.notes

    def _add_metrics_slide(self, config: SlideConfig):
        """Add a metrics dashboard slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        self._add_slide_header(slide, config.title)

        metrics = config.content or {}
        items = list(metrics.items())

        # Create 2x3 grid of metric boxes
        cols = 3
        rows = (len(items) + cols - 1) // cols

        box_width = 4
        box_height = 1.8
        start_x = 0.5
        start_y = 1.5
        gap = 0.2

        for i, (label, value) in enumerate(items):
            row = i // cols
            col = i % cols

            x = start_x + col * (box_width + gap)
            y = start_y + row * (box_height + gap)

            # Box background
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(box_width), Inches(box_height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
            shape.line.color.rgb = RGBColor(220, 220, 220)

            # Accent bar
            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y),
                Inches(0.1), Inches(box_height)
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor(*self.config.primary_color)
            accent.line.fill.background()

            # Label
            label_box = slide.shapes.add_textbox(
                Inches(x + 0.3), Inches(y + 0.3),
                Inches(box_width - 0.4), Inches(0.4)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(100, 100, 100)

            # Value
            value_box = slide.shapes.add_textbox(
                Inches(x + 0.3), Inches(y + 0.8),
                Inches(box_width - 0.4), Inches(0.8)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(value)
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*self.config.primary_color)

    def _add_table_slide(self, config: SlideConfig):
        """Add a table slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        self._add_slide_header(slide, config.title)

        data = config.content or {}
        headers = data.get("headers", [])
        rows = data.get("rows", [])

        if not headers and not rows:
            return

        all_rows = [headers] + rows if headers else rows
        num_cols = len(all_rows[0]) if all_rows else 0
        num_rows = len(all_rows)

        if num_cols == 0:
            return

        # Calculate table dimensions
        table_width = min(12.333, num_cols * 2)
        col_width = table_width / num_cols

        table = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(0.5), Inches(1.5),
            Inches(table_width), Inches(min(5, num_rows * 0.5))
        ).table

        # Set column widths
        for i in range(num_cols):
            table.columns[i].width = Inches(col_width)

        # Fill data
        for row_idx, row_data in enumerate(all_rows):
            for col_idx, cell_value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_value)

                # Style
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)
                para.alignment = PP_ALIGN.CENTER

                if row_idx == 0 and headers:
                    # Header row
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*self.config.primary_color)
                    para.font.color.rgb = RGBColor(255, 255, 255)
                    para.font.bold = True
                else:
                    # Data rows - alternate colors
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(250, 250, 250)

    def _add_chart_slide(self, config: SlideConfig):
        """Add a chart slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        self._add_slide_header(slide, config.title)

        chart_data = config.content or {}
        chart_type = chart_data.get("type", "bar")
        categories = chart_data.get("categories", [])
        series_list = chart_data.get("series", [])

        if not categories or not series_list:
            return

        data = CategoryChartData()
        data.categories = categories

        for series in series_list:
            data.add_series(series.get("name", "Series"), series.get("values", []))

        # Determine chart type
        if chart_type == "line":
            xl_chart_type = XL_CHART_TYPE.LINE
        elif chart_type == "bar":
            xl_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        elif chart_type == "stacked":
            xl_chart_type = XL_CHART_TYPE.COLUMN_STACKED
        else:
            xl_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED

        chart = slide.shapes.add_chart(
            xl_chart_type, Inches(1), Inches(1.5),
            Inches(11), Inches(5.5),
            data
        ).chart

        # Style the chart
        chart.has_legend = len(series_list) > 1
        if chart.has_legend:
            chart.legend.include_in_layout = False

    def _add_comparison_slide(self, config: SlideConfig):
        """Add a side-by-side comparison slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        self._add_slide_header(slide, config.title)

        comparisons = config.content or []
        if len(comparisons) < 2:
            return

        # Left side
        left_title = comparisons[0].get("title", "Option A")
        left_items = comparisons[0].get("items", [])

        left_header = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(6), Inches(0.5)
        )
        tf = left_header.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.config.primary_color)

        left_content = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.2),
            Inches(6), Inches(4.5)
        )
        tf = left_content.text_frame
        for i, item in enumerate(left_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.space_after = Pt(8)

        # Divider
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.5), Inches(1.5),
            Inches(0.02), Inches(5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(200, 200, 200)
        line.line.fill.background()

        # Right side
        right_title = comparisons[1].get("title", "Option B")
        right_items = comparisons[1].get("items", [])

        right_header = slide.shapes.add_textbox(
            Inches(6.833), Inches(1.5),
            Inches(6), Inches(0.5)
        )
        tf = right_header.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.config.accent_color)

        right_content = slide.shapes.add_textbox(
            Inches(6.833), Inches(2.2),
            Inches(6), Inches(4.5)
        )
        tf = right_content.text_frame
        for i, item in enumerate(right_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.space_after = Pt(8)

    def _add_slide_header(self, slide, title: str):
        """Add a header to a slide."""
        # Header bar
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(1.2)
        )
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = RGBColor(*self.config.primary_color)
        header_bar.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.35),
            Inches(12.333), Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Footer
        footer = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.1),
            Inches(3), Inches(0.3)
        )
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = self.config.footer_text
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(150, 150, 150)

        # Page number placeholder (bottom right)
        page_num = slide.shapes.add_textbox(
            Inches(12), Inches(7.1),
            Inches(0.833), Inches(0.3)
        )
        tf = page_num.text_frame
        p = tf.paragraphs[0]
        p.text = str(len(self.prs.slides))
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(150, 150, 150)

    # ===== Specialized Presentation Generators =====

    def generate_lbo_presentation(self, outputs: Dict[str, Any]) -> bytes:
        """Generate an LBO analysis presentation."""
        slides = []

        # Executive Summary
        slides.append(SlideConfig(
            title="Executive Summary",
            slide_type="metrics",
            content={
                "IRR": f"{outputs.get('irr', 0):.1%}",
                "MOIC": f"{outputs.get('moic', 0):.2f}x",
                "Entry EV/EBITDA": f"{outputs.get('entry_ev_ebitda', 0):.1f}x",
                "Exit Equity Value": f"${outputs.get('exit_equity_value', 0):,.0f}",
                "Total Invested": f"${outputs.get('total_equity_invested', 0):,.0f}",
                "Holding Period": f"{outputs.get('exit_year', 5)} Years",
            }
        ))

        # Sources & Uses
        sources = outputs.get('sources', {})
        uses = outputs.get('uses', {})

        su_headers = ["Sources", "Amount ($M)", "Uses", "Amount ($M)"]
        su_rows = []

        source_items = [(k, v) for k, v in sources.items() if v > 0]
        use_items = [(k, v) for k, v in uses.items() if v > 0]
        max_len = max(len(source_items), len(use_items))

        for i in range(max_len):
            row = []
            if i < len(source_items):
                row.extend([source_items[i][0], f"{source_items[i][1]:,.1f}"])
            else:
                row.extend(["", ""])
            if i < len(use_items):
                row.extend([use_items[i][0], f"{use_items[i][1]:,.1f}"])
            else:
                row.extend(["", ""])
            su_rows.append(row)

        slides.append(SlideConfig(
            title="Sources & Uses of Funds",
            slide_type="table",
            content={"headers": su_headers, "rows": su_rows}
        ))

        # Projections Chart
        years = outputs.get('years', [])
        revenues = outputs.get('revenues', [])
        ebitda = outputs.get('ebitda', [])

        if years and revenues:
            slides.append(SlideConfig(
                title="Revenue & EBITDA Projections",
                slide_type="chart",
                content={
                    "type": "bar",
                    "categories": [f"Year {y}" for y in years],
                    "series": [
                        {"name": "Revenue", "values": revenues},
                        {"name": "EBITDA", "values": ebitda},
                    ]
                }
            ))

        # Returns Sensitivity (if available)
        slides.append(SlideConfig(
            title="Key Investment Considerations",
            slide_type="content",
            content=[
                f"Target IRR of {outputs.get('irr', 0):.1%} achieved through operational improvements and debt paydown",
                f"Entry multiple of {outputs.get('entry_ev_ebitda', 0):.1f}x provides downside protection",
                f"Equity contribution of {outputs.get('equity_contribution_percent', 0):.0%} aligned with market standards",
                "Strong free cash flow generation supports aggressive debt amortization",
                "Multiple expansion potential through margin improvement initiatives",
            ]
        ))

        self.config.title = "LBO Analysis"
        self.config.subtitle = "Investment Committee Presentation"
        return self.generate_presentation(slides)

    def generate_valuation_presentation(self, dcf: Dict, comps: Dict = None, precedents: Dict = None) -> bytes:
        """Generate a valuation summary presentation."""
        slides = []

        # Valuation Summary
        metrics = {
            "DCF Value": f"${dcf.get('equity_value_per_share', 0):,.2f}",
            "Enterprise Value": f"${dcf.get('enterprise_value', 0):,.0f}M",
            "WACC": f"{dcf.get('wacc', 0):.1%}",
            "Terminal Growth": f"{dcf.get('terminal_growth_rate', 0):.1%}",
        }

        if comps:
            metrics["Comps Median"] = f"${comps.get('median_value', 0):,.2f}"
        if precedents:
            metrics["Precedents Median"] = f"${precedents.get('median_value', 0):,.2f}"

        slides.append(SlideConfig(
            title="Valuation Summary",
            slide_type="metrics",
            content=metrics
        ))

        # DCF Build
        if dcf.get('projected_fcf'):
            slides.append(SlideConfig(
                title="DCF Cash Flow Projections",
                slide_type="chart",
                content={
                    "type": "bar",
                    "categories": [f"Year {i+1}" for i in range(len(dcf['projected_fcf']))],
                    "series": [{"name": "Free Cash Flow", "values": dcf['projected_fcf']}]
                }
            ))

        # Comparable Companies
        if comps and comps.get('companies'):
            comp_rows = []
            for comp in comps['companies'][:8]:
                comp_rows.append([
                    comp.get('name', ''),
                    f"{comp.get('ev_ebitda', 0):.1f}x",
                    f"{comp.get('ev_revenue', 0):.1f}x",
                    f"{comp.get('pe', 0):.1f}x",
                ])

            slides.append(SlideConfig(
                title="Trading Comparables",
                slide_type="table",
                content={
                    "headers": ["Company", "EV/EBITDA", "EV/Revenue", "P/E"],
                    "rows": comp_rows
                }
            ))

        self.config.title = "Valuation Analysis"
        return self.generate_presentation(slides)

    def generate_scenario_comparison(self, scenarios: Dict[str, Dict]) -> bytes:
        """Generate a scenario comparison presentation."""
        slides = []

        scenario_names = list(scenarios.keys())

        # Summary comparison
        if len(scenario_names) >= 2:
            base = scenarios.get(scenario_names[0], {})
            alt = scenarios.get(scenario_names[1], {})

            slides.append(SlideConfig(
                title="Scenario Comparison",
                slide_type="comparison",
                content=[
                    {
                        "title": scenario_names[0],
                        "items": [
                            f"IRR: {base.get('irr', 0):.1%}",
                            f"MOIC: {base.get('moic', 0):.2f}x",
                            f"Exit Value: ${base.get('exit_equity_value', 0):,.0f}",
                        ]
                    },
                    {
                        "title": scenario_names[1],
                        "items": [
                            f"IRR: {alt.get('irr', 0):.1%}",
                            f"MOIC: {alt.get('moic', 0):.2f}x",
                            f"Exit Value: ${alt.get('exit_equity_value', 0):,.0f}",
                        ]
                    }
                ]
            ))

        # All scenarios table
        if scenarios:
            headers = ["Metric"] + scenario_names
            rows = [
                ["IRR"] + [f"{s.get('irr', 0):.1%}" for s in scenarios.values()],
                ["MOIC"] + [f"{s.get('moic', 0):.2f}x" for s in scenarios.values()],
                ["Exit Value"] + [f"${s.get('exit_equity_value', 0):,.0f}" for s in scenarios.values()],
            ]

            slides.append(SlideConfig(
                title="Scenario Summary",
                slide_type="table",
                content={"headers": headers, "rows": rows}
            ))

        self.config.title = "Scenario Analysis"
        return self.generate_presentation(slides)
